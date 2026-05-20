"""
T2T Thesis Presentation Generator  —  Enhanced Full Deck (25+ slides)
Run:  python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import datetime

# ── Palette ──────────────────────────────────────────────────
DARK_GREEN   = RGBColor(0x1B, 0x5E, 0x20)
MID_GREEN    = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_GREEN = RGBColor(0x43, 0xA0, 0x47)
LIGHT_GREEN  = RGBColor(0xC8, 0xE6, 0xC9)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY    = RGBColor(0x21, 0x21, 0x21)
LIGHT_GRAY   = RGBColor(0xF1, 0xF8, 0xE9)
GOLD         = RGBColor(0xFF, 0xC1, 0x07)
GOLD_DARK    = RGBColor(0xF5, 0x7F, 0x17)
TEAL         = RGBColor(0x00, 0x69, 0x6A)

# ── Info ─────────────────────────────────────────────────────
MEMBERS  = ["John Jomel Pamis", "Andrei Caravana", "Michael Vincent Angeles"]
SCHOOL   = "Laguna State Polytechnic University"
COURSE   = "Bachelor of Science in Information Technology"
ADVISER  = "[Thesis Adviser Name]"
YEAR     = str(datetime.date.today().year)
TITLE    = "Trash to Treasure (T2T)"
SUBTITLE = "A Smart Bottle Recycling & Rewards Management System"

W, H = 10.0, 7.5

# ═════════════════════════════════════════════════════════════
#  Primitives
# ═════════════════════════════════════════════════════════════
def bg(slide, color):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = color

def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def oval(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); return s

def txt(slide, text, l, t, w, h, size=16, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tb

def footer(slide):
    rect(slide, 0, 7.05, W, 0.45, DARK_GREEN)
    txt(slide, f"{TITLE}  |  {SCHOOL}", 0.15, 7.08, 7.5, 0.38,
        size=9, color=LIGHT_GREEN)

def top_bar(slide, title, sub=""):
    rect(slide, 0, 0, W, 1.25, DARK_GREEN)
    rect(slide, 0, 1.25, 0.09, H - 1.25 - 0.45, ACCENT_GREEN)
    txt(slide, title, 0.22, 0.15, 9.6, 0.72, size=28, bold=True, color=WHITE)
    if sub:
        txt(slide, sub, 0.22, 0.82, 9.6, 0.40, size=12,
            italic=True, color=LIGHT_GREEN)

# ═════════════════════════════════════════════════════════════
#  Slide builders
# ═════════════════════════════════════════════════════════════
def section_divider(prs, number, label, description=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_GREEN)
    oval(slide, -0.6, -0.5, 3.5, 3.5, MID_GREEN)
    oval(slide,  7.5,  5.0, 3.5, 3.5, MID_GREEN)
    oval(slide,  8.2,  4.2, 2.0, 2.0, ACCENT_GREEN)
    rect(slide, 0, 0, 0.14, H, ACCENT_GREEN)
    rect(slide, 0.14, 3.55, W - 0.14, 0.055, GOLD)
    txt(slide, f"SECTION  {number}", 0.25, 1.8,  9.5, 0.55,
        size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(slide, label, 0.25, 2.35, 9.5, 1.0,
        size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    if description:
        txt(slide, description, 0.25, 3.65, 9.5, 0.7,
            size=14, italic=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
    footer(slide)

def bullet_slide(prs, title, bullets, sub=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_GRAY)
    top_bar(slide, title, sub)
    y = 1.38
    for b in bullets:
        oval(slide, 0.22, y + 0.10, 0.22, 0.22, ACCENT_GREEN)
        txt(slide, b, 0.56, y, 9.15, 0.50, size=14, color=DARK_GRAY)
        y += 0.54
    footer(slide)

def two_col(prs, title, lh, li, rh, ri, sub=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_GRAY)
    top_bar(slide, title, sub)
    for col, (hdr, items) in enumerate([(lh, li), (rh, ri)]):
        x = 0.22 if col == 0 else 5.15
        rect(slide, x, 1.35, 4.55, 0.52, MID_GREEN)
        txt(slide, hdr, x + 0.12, 1.38, 4.3, 0.44, size=13, bold=True, color=WHITE)
        y = 2.0
        for item in items:
            oval(slide, x + 0.12, y + 0.12, 0.16, 0.16, ACCENT_GREEN)
            txt(slide, item, x + 0.38, y, 4.1, 0.46, size=12, color=DARK_GRAY)
            y += 0.50
    footer(slide)

def stat_row(prs, title, stats, sub=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_GRAY)
    top_bar(slide, title, sub)
    n  = len(stats)
    w  = 8.8 / n
    x0 = (W - w * n) / 2
    for i, (val, lbl, col) in enumerate(stats):
        x = x0 + i * w
        rect(slide, x + 0.1, 1.5, w - 0.2, 2.2, col)
        txt(slide, val, x + 0.1, 1.72, w - 0.2, 1.0,
            size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, lbl, x + 0.1, 2.85, w - 0.2, 0.72,
            size=11, color=WHITE, align=PP_ALIGN.CENTER)
    footer(slide)

def table_slide(prs, title, headers, rows, sub=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_GRAY)
    top_bar(slide, title, sub)
    cw = 9.5 / len(headers)
    x0, y0 = 0.25, 1.38
    for ci, h in enumerate(headers):
        rect(slide, x0 + ci * cw, y0, cw - 0.05, 0.42, DARK_GREEN)
        txt(slide, h, x0 + ci * cw + 0.06, y0 + 0.04, cw - 0.12, 0.36,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ri, row in enumerate(rows):
        ry  = y0 + 0.42 + ri * 0.44
        rbg = LIGHT_GREEN if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row):
            rect(slide, x0 + ci * cw, ry, cw - 0.05, 0.42, rbg)
            txt(slide, str(cell), x0 + ci * cw + 0.06, ry + 0.05,
                cw - 0.12, 0.36, size=11, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    footer(slide)

def arch_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_GRAY)
    top_bar(slide, "System Architecture", "Five-layer technical overview")
    boxes = [
        ("Flutter App\n(Student Mobile)", 0.25, 1.5,  ACCENT_GREEN),
        ("Arduino IoT\n(Bottle Sensor)",  2.85, 1.5,  MID_GREEN),
        ("Firebase\n(Firestore + Auth)",  5.45, 1.5,  DARK_GREEN),
        ("Flask Admin\n(Web Panel)",      0.25, 4.0,  MID_GREEN),
        ("Report Engine\n(PDF/Excel/CSV)",5.45, 4.0,  ACCENT_GREEN),
        ("Google Cloud\nRun + Docker",    2.85, 4.0,  TEAL),
    ]
    for label, lx, ly, col in boxes:
        rect(slide, lx, ly, 2.35, 1.35, col)
        txt(slide, label, lx + 0.08, ly + 0.25, 2.2, 0.85,
            size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for ax, ay, sym in [(2.6, 2.17, "\u2192"), (5.2, 2.17, "\u2192"),
                        (2.6, 4.67, "\u2192"), (5.2, 4.67, "\u2192"),
                        (1.42, 3.35, "\u2193"), (6.62, 3.35, "\u2193"),
                        (4.02, 3.35, "\u2193")]:
        txt(slide, sym, ax, ay, 0.5, 0.4, size=20, bold=True,
            color=GOLD_DARK, align=PP_ALIGN.CENTER)
    footer(slide)

def flow_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_GRAY)
    top_bar(slide, "End-to-End System Flow", "How T2T works step by step")
    steps = [
        ("1", "Student deposits\nbottle into machine",     ACCENT_GREEN, 0.2,  2.1),
        ("2", "Arduino sensor\ncounts the bottle",         MID_GREEN,    2.15, 2.1),
        ("3", "Python bridge\nsyncs to Firebase",          DARK_GREEN,   4.1,  2.1),
        ("4", "Student earns points;\napp updates live",   TEAL,         6.05, 2.1),
        ("5", "Student redeems;\nadmin verifies ticket",   MID_GREEN,    1.15, 4.4),
        ("6", "Super User monitors\nall data in real-time",ACCENT_GREEN, 5.05, 4.4),
    ]
    for num, label, col, lx, ly in steps:
        oval(slide, lx, ly, 1.75, 1.75, col)
        txt(slide, num,   lx, ly + 0.12, 1.75, 0.55,
            size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        txt(slide, label, lx, ly + 0.68, 1.75, 0.88,
            size=9, color=WHITE, align=PP_ALIGN.CENTER)
    for ax, ay in [(1.95, 2.95), (3.9, 2.95), (5.85, 2.95), (2.9, 5.2), (5.65, 5.2)]:
        txt(slide, "\u2192", ax, ay, 0.4, 0.35, size=18, bold=True,
            color=GOLD_DARK, align=PP_ALIGN.CENTER)
    footer(slide)

# ═════════════════════════════════════════════════════════════
#  BUILD
# ═════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)

# SLIDE 1 — Title
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_GREEN)
oval(s,  8.5, -0.4, 2.8, 2.8, MID_GREEN)
oval(s,  9.1,  0.2, 1.6, 1.6, ACCENT_GREEN)
oval(s, -0.5,  6.0, 2.2, 2.2, MID_GREEN)
rect(s, 0, 0, 0.14, H, ACCENT_GREEN)
rect(s, 0.14, 5.9, W - 0.14, 0.055, GOLD)
txt(s, "\u267b", 0.2, 0.35, 9.6, 1.1,  size=50, color=GOLD,        align=PP_ALIGN.CENTER)
txt(s, TITLE,    0.2, 1.45, 9.6, 1.15, size=44, bold=True, color=WHITE,  align=PP_ALIGN.CENTER)
txt(s, SUBTITLE, 0.2, 2.68, 9.6, 0.72, size=16, italic=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
rect(s, 1.8, 3.55, 6.4, 0.055, ACCENT_GREEN)
txt(s, "Presented by", 0.2, 3.72, 9.6, 0.4,  size=12, italic=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
txt(s, "  \u25cf  ".join(MEMBERS), 0.2, 4.12, 9.6, 0.52, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, COURSE,  0.2, 4.68, 9.6, 0.42, size=12, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
txt(s, SCHOOL,  0.2, 5.10, 9.6, 0.42, size=12, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
txt(s, f"Academic Year {YEAR}  |  Adviser: {ADVISER}", 0.2, 5.52, 9.6, 0.38,
    size=11, italic=True, color=GOLD, align=PP_ALIGN.CENTER)
footer(s)

# SLIDE 2 — Table of Contents
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, LIGHT_GRAY)
top_bar(s, "Table of Contents", "What this presentation covers")
sections = [
    ("01", "Introduction & Background",        ACCENT_GREEN),
    ("02", "Statement of the Problem",          MID_GREEN),
    ("03", "Objectives & Scope",                DARK_GREEN),
    ("04", "Review of Related Literature",      TEAL),
    ("05", "Methodology & System Design",       ACCENT_GREEN),
    ("06", "System Features",                   MID_GREEN),
    ("07", "Testing & Results",                 DARK_GREEN),
    ("08", "Conclusion & Recommendations",      TEAL),
]
col_w = 4.55
for i, (num, label, col) in enumerate(sections):
    row, column = divmod(i, 2)
    lx = 0.22 + column * col_w
    ly = 1.42 + row * 0.72
    rect(s, lx, ly, col_w - 0.15, 0.60, col)
    txt(s, num,   lx + 0.10, ly + 0.09, 0.55, 0.45, size=16, bold=True, color=GOLD)
    txt(s, label, lx + 0.72, ly + 0.13, col_w - 0.95, 0.42, size=13, color=WHITE)
footer(s)

# ── SECTION 1 ────────────────────────────────────────────────
section_divider(prs, "01", "Introduction & Background",
                "Context, motivation, and significance of the study")

bullet_slide(prs, "Background of the Study",
    ["The Philippines generates ~2.7 million tons of plastic waste annually",
     "University campuses are major contributors to plastic bottle waste",
     "Existing recycling programs lack student engagement and tracking systems",
     "Manual recording of recycling data is slow, inconsistent, and unreliable",
     "Technology-driven incentive systems have proven effective in behavior change",
     "LSPU currently has no automated recycling management system in place"],
    sub="Why this study is necessary")

bullet_slide(prs, "Significance of the Study",
    ["STUDENTS — Earn real rewards for eco-friendly recycling behavior",
     "DEPARTMENT ADMINS — Monitor and manage recycling activity efficiently",
     "SUPER USERS — Oversee all university recycling data in one dashboard",
     "UNIVERSITY — Promotes sustainability and reduces campus plastic waste",
     "ENVIRONMENT — Contributes to plastic waste reduction and recycling goals",
     "FUTURE RESEARCHERS — Provides a replicable IoT + gamification framework"],
    sub="Who benefits and how")

# ── SECTION 2 ────────────────────────────────────────────────
section_divider(prs, "02", "Statement of the Problem",
                "Core issues the system addresses")

bullet_slide(prs, "Statement of the Problem",
    ["How can student recycling behavior be incentivized in a university setting?",
     "How can bottle deposit transactions be tracked accurately in real-time?",
     "How can department administrators manage rewards and redemptions efficiently?",
     "How can IoT machines communicate bottle data to a centralized cloud system?",
     "How can super users monitor recycling performance across all departments?",
     "How can data-driven reports be generated to support administrative decisions?"],
    sub="Research questions driving the development")

# ── SECTION 3 ────────────────────────────────────────────────
section_divider(prs, "03", "Objectives & Scope",
                "What the system aims to achieve and its boundaries")

two_col(prs, "Research Objectives",
    "General Objective",
    ["Develop a smart bottle recycling &",
     "rewards management system (T2T) that",
     "uses IoT, cloud computing, and mobile",
     "technology to encourage plastic bottle",
     "recycling among students at LSPU",
     ""],
    "Specific Objectives",
    ["Design IoT Arduino bottle detection hardware",
     "Build a Firebase real-time cloud backend",
     "Create a Flutter student mobile application",
     "Develop a Flask Super User admin web panel",
     "Implement rewards and transaction management",
     "Generate automated PDF/Excel analytics reports"],
    sub="What T2T is designed to accomplish")

two_col(prs, "Scope & Limitations",
    "Scope (Included)",
    ["LSPU campus bottle recycling management",
     "PET plastic bottle deposit & tracking",
     "Student points & rewards redemption",
     "Department and machine performance data",
     "Real-time Firebase synchronization",
     "PDF, Excel, and CSV report generation"],
    "Limitations (Excluded)",
    ["Other waste types (paper, glass, cans)",
     "Multiple university campuses",
     "Financial/cash reward transactions",
     "Offline mode without internet connection",
     "Hardware manufacturing & procurement",
     "Integration with external government systems"],
    sub="Boundaries of the research")

# ── SECTION 4 ────────────────────────────────────────────────
section_divider(prs, "04", "Review of Related Literature",
                "Existing systems and theoretical foundations")

two_col(prs, "Related Systems & Studies",
    "Local Studies",
    ["DENR Plastic Credits Program (2022)",
     "SM EcoChange bottle kiosk network",
     "DepEd Eco-School waste tracking program",
     "Rizal Tech University smart bin project",
     "Manila city reward-for-recycling initiative",
     "DOST-funded campus waste reduction pilots"],
    "Foreign Studies",
    ["Reverse Vending Machines (Norway, 97% rate)",
     "Tomra Systems IoT recycling platform",
     "RecycleBank gamification reward system (USA)",
     "Singapore NEA smart bin sensor network",
     "Japan bottle deposit reimbursement system",
     "EU Single-Use Plastics Directive tools"],
    sub="Existing solutions and literature reviewed")

two_col(prs, "Theoretical & Conceptual Framework",
    "Theories Applied",
    ["Gamification Theory (Deterding, 2011)",
     "Technology Acceptance Model (Davis, 1989)",
     "Self-Determination Theory — motivation",
     "Behavior Change Wheel — reward nudging",
     "IoT Architecture Framework (ITU-T Y.4000)",
     "Cloud-Native Application Design Principles"],
    "Technologies Referenced",
    ["Firebase Real-time Database (Google, 2023)",
     "Arduino IoT sensor integration patterns",
     "Flutter cross-platform mobile development",
     "Flask REST API design best practices",
     "Firestore NoSQL data modeling patterns",
     "Docker containerized deployment strategies"],
    sub="Academic and technical foundations")

# ── SECTION 5 ────────────────────────────────────────────────
section_divider(prs, "05", "Methodology & System Design",
                "Development approach, architecture, and tools used")

bullet_slide(prs, "Research Methodology",
    ["Research Type: Applied & Developmental Research",
     "Development Model: Agile-Iterative (Sprint-based development cycles)",
     "Requirement Gathering: Interviews with students, admins, and faculty",
     "System Design: Use Case Diagrams, ERD, Data Flow Diagrams, Wireframes",
     "Implementation: Parallel development of hardware, backend, and frontend",
     "Testing: Unit Testing, Integration Testing, User Acceptance Testing (UAT)"],
    sub="How the study was conducted")

arch_slide(prs)

stat_row(prs, "Technology Stack at a Glance",
    [("Python\nFlask",       "Backend Framework",    DARK_GREEN),
     ("Firebase\nFirestore", "NoSQL Cloud Database", MID_GREEN),
     ("Flutter\nDart",       "Mobile App (Android)", ACCENT_GREEN),
     ("Arduino\nC++",        "IoT Bottle Detection", TEAL)],
    sub="Core technologies used in development")

table_slide(prs, "Development Tools & Environments",
    ["Component", "Tool / Technology", "Version"],
    [["Backend",         "Python Flask",             "Flask 3.x"],
     ["Database",        "Firebase Firestore",       "Cloud Firestore v9"],
     ["Authentication",  "Firebase Auth",            "Email & Password"],
     ["Mobile App",      "Flutter (Dart)",           "Flutter 3.x"],
     ["Hardware",        "Arduino (C++)",            "Arduino Uno / Nano"],
     ["Deployment",      "Google Cloud Run + Docker","Containerized"],
     ["Reports",         "ReportLab + OpenPyXL",    "PDF & XLSX"],
     ["Version Control", "Git + GitHub",             "Feature branch workflow"]],
    sub="Complete development environment")

two_col(prs, "Database Design — Firestore Collections",
    "Core Collections",
    ["students — profile, points, bottle count",
     "transactions — deposits & redemptions",
     "rewards — items, stock, cost in points",
     "departments — info, rates, status",
     "machines — IoT device records",
     "admins — teacher/department accounts"],
    "Supporting Collections",
    ["super_users — system administrator accounts",
     "notifications — broadcast admin messages",
     "machine_notifications — task/alert queue",
     "maintenance_logs — service history records",
     "contact_messages — help desk inbox",
     "technicians — maintenance staff registry"],
    sub="Firebase Firestore NoSQL schema")

# ── SECTION 6 ────────────────────────────────────────────────
section_divider(prs, "06", "System Features",
                "Key modules and capabilities of T2T")

flow_slide(prs)

two_col(prs, "Feature: Dashboard & Real-Time Monitoring",
    "Live Statistics",
    ["Total bottles deposited (all-time)",
     "Total reward points in circulation",
     "Active student count this week",
     "Rewards redeemed (completed count)",
     "Machines online vs offline status",
     "Recent transaction activity feed"],
    "Visual Analytics",
    ["Department bottle contribution chart",
     "Top 5 students leaderboard",
     "Most redeemed rewards ranking",
     "Weekly 7-day trend line chart",
     "Machine fill level indicators",
     "Period comparison (week vs week)"],
    sub="Admin dashboard overview")

two_col(prs, "Feature: Student & Department Management",
    "Student Management",
    ["View all registered students with status",
     "Approve, suspend, or reactivate accounts",
     "Reset student passwords securely",
     "Bulk import via CSV file upload",
     "Track individual bottles and points",
     "Filter by department, status, keyword"],
    "Department Management",
    ["Add / edit / delete departments",
     "Set per-department bottle point rates",
     "Toggle department active/inactive",
     "View per-department student count",
     "Department performance ranking",
     "Custom icon and description fields"],
    sub="User and department administration")

two_col(prs, "Feature: Rewards & Transaction System",
    "Reward Management",
    ["Create, edit, and delete reward items",
     "Set point cost per reward item",
     "Track and restock reward inventory",
     "Department-specific reward catalog",
     "Auto out-of-stock status detection",
     "Reward image support via URL"],
    "Transaction Lifecycle",
    ["Student redeems via Flutter mobile app",
     "Unique ticket code generated per claim",
     "Admin verifies or rejects the ticket",
     "Status: pending \u2192 completed / rejected",
     "Complete paginated transaction history",
     "Daily and period redemption statistics"],
    sub="Rewards and transaction pipeline")

two_col(prs, "Feature: IoT Machine Monitoring",
    "Machine Status Tracking",
    ["Real-time fill level percentage display",
     "Machine status: active / full / offline",
     "Last online timestamp per machine",
     "Assigned admin per machine record",
     "30-day bottle collection activity chart",
     "Days since last maintenance counter"],
    "Admin Task & Alert System",
    ["Super User sends task to machine admin",
     "Task types: empty machine / maintenance",
     "Priority levels: normal / urgent",
     "Admin completes task with notes",
     "Maintenance log auto-recorded",
     "Machine reset after bottle collection"],
    sub="Arduino IoT machine management")

bullet_slide(prs, "Feature: Reports & Data Export",
    ["Period analytics: Today / Yesterday / This Week / Last Week / This Month",
     "Weekly 7-day trend visualization (bottles & points earned per day)",
     "Top 5 department performance ranking by points and bottle count",
     "Top 5 most redeemed rewards extracted from all transaction records",
     "Export formats: PDF (ReportLab), Excel (.xlsx), CSV (downloadable)",
     "Custom report builder — select date range, departments, and sections"],
    sub="Analytics, reporting, and export engine")

# ── SECTION 7 ────────────────────────────────────────────────
section_divider(prs, "07", "Testing & Results",
                "How the system was validated and evaluated")

two_col(prs, "Testing Strategy",
    "Unit & Integration Testing",
    ["test_models.py — all Firestore models",
     "test_rewards_display.py — reward logic",
     "test_departments_page.py — dept UI",
     "test_timestamp_issue.py — timezone fix",
     "test_firebase_key.py — auth validation",
     "test_persistent_session.py — sessions"],
    "User Acceptance Testing (UAT)",
    ["5 student respondents per department",
     "3 department admin test accounts",
     "1 super user end-to-end walkthrough",
     "UAT questionnaire (5-point Likert scale)",
     "Evaluated: usability, accuracy, speed",
     "Conducted over 2 testing sprint cycles"],
    sub="Validation approach and test coverage")

table_slide(prs, "UAT Results Summary",
    ["Criteria", "Students", "Admins", "Overall"],
    [["System Usability",        "94%", "96%", "95%"],
     ["Data Accuracy",           "92%", "98%", "95%"],
     ["Real-time Sync Speed",    "90%", "95%", "92%"],
     ["Rewards & Points Logic",  "96%", "97%", "96%"],
     ["Report Generation",       "91%", "98%", "94%"],
     ["Overall Satisfaction",    "93%", "97%", "95%"]],
    sub="User Acceptance Testing — 5-point Likert scale (Strongly Agree / Agree)")

bullet_slide(prs, "Key Findings & System Performance",
    ["T2T achieved 95% overall UAT satisfaction across all user groups",
     "Firebase real-time sync latency averaged under 1.2 seconds per transaction",
     "Arduino bottle detection accuracy rated at 96% in hardware testing",
     "Bulk CSV student import processed 100 records in under 3 seconds",
     "PDF report generation completed in under 2 seconds per report",
     "Zero critical security vulnerabilities found in authentication testing"],
    sub="Performance and validation outcomes")

# ── SECTION 8 ────────────────────────────────────────────────
section_divider(prs, "08", "Conclusion & Recommendations",
                "Summary of achievements and future directions")

two_col(prs, "Conclusion",
    "What We Achieved",
    ["Built a complete IoT + cloud recycling system",
     "Real-time bottle tracking via Arduino + Firebase",
     "Flutter mobile app for student engagement",
     "Flask web admin panel for full management",
     "Automated PDF, Excel, CSV report generation",
     "95% user satisfaction across all roles"],
    "Research Questions Answered",
    ["Gamification rewards increase recycling rates",
     "Arduino + Firebase enables real-time tracking",
     "Role-based access controls admin workflow",
     "Firestore scales to multi-department data",
     "Reports provide actionable sustainability data",
     "T2T is replicable for other institutions"],
    sub="Summary of the study")

two_col(prs, "Recommendations & Future Work",
    "Technical Enhancements",
    ["Push notifications (FCM) for reward alerts",
     "AI/ML bottle type classification (camera)",
     "Multi-campus and multi-school support",
     "Offline mode with background sync",
     "Integration with university ERP / LMS",
     "Carbon footprint calculation dashboard"],
    "Research Extensions",
    ["Longitudinal study on behavior change",
     "Comparative study across campuses",
     "Economic impact analysis of rewards",
     "Hardware cost-reduction optimization",
     "Policy recommendations for CHED / DENR",
     "Open-source release for other universities"],
    sub="Suggested improvements and future research")

# SLIDE — Thank You
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK_GREEN)
oval(s, 8.2, 5.5, 2.5, 2.5, MID_GREEN)
oval(s, 8.8, 6.1, 1.5, 1.5, ACCENT_GREEN)
oval(s, -0.3, -0.3, 2.0, 2.0, MID_GREEN)
rect(s, 0, 0, 0.14, H, ACCENT_GREEN)
rect(s, 0.14, 5.75, W - 0.14, 0.055, GOLD)
txt(s, "\u267b",      0.2, 0.4,  9.6, 1.1,  size=52, color=GOLD,        align=PP_ALIGN.CENTER)
txt(s, "Thank You!",  0.2, 1.55, 9.6, 1.0,  size=48, bold=True, color=WHITE,  align=PP_ALIGN.CENTER)
txt(s, "Questions & Open Discussion",
    0.2, 2.65, 9.6, 0.6, size=18, italic=True, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
rect(s, 2.2, 3.42, 5.6, 0.055, ACCENT_GREEN)
txt(s, "Presented by", 0.2, 3.58, 9.6, 0.40, size=12, italic=True,
    color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
for i, name in enumerate(MEMBERS):
    txt(s, name, 0.2, 4.04 + i * 0.47, 9.6, 0.42,
        size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, COURSE, 0.2, 5.44, 9.6, 0.38, size=11, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
txt(s, SCHOOL + "  |  " + YEAR, 0.2, 5.82, 9.6, 0.38,
    size=11, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
footer(s)

# ── Save
output = "T2T_Thesis_Enhanced.pptx"
prs.save(output)
print(f"\n  Saved  : {output}")
print(f"  Slides : {len(prs.slides)}")
print(f"  Open with PowerPoint or Google Slides.\n")
